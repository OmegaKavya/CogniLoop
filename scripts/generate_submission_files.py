import os
import sys
import re
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether, HRFlowable
from reportlab.pdfgen import canvas
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

os.makedirs("Submission", exist_ok=True)

class NumberedCanvas(canvas.Canvas):
    def __init__(self, *args, **kwargs):
        super(NumberedCanvas, self).__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_decorations(num_pages)
            super(NumberedCanvas, self).showPage()
        super(NumberedCanvas, self).save()

    def draw_page_decorations(self, page_count):
        self.saveState()
        self.setFont("Helvetica", 8)
        self.setFillColor(colors.HexColor("#71717a"))
        
        # Header (pages > 1)
        if self._pageNumber > 1:
            self.drawString(54, 750, "CogniLoop Official Research & Technical Submission")
            self.setStrokeColor(colors.HexColor("#e2e8f0"))
            self.setLineWidth(0.5)
            self.line(54, 742, 558, 742)
            
        # Footer
        self.setStrokeColor(colors.HexColor("#e2e8f0"))
        self.setLineWidth(0.5)
        self.line(54, 45, 558, 45)
        
        self.drawString(54, 32, "Author: Kavya Aggarwal | CogniLoop v1.0")
        page_str = f"Page {self._pageNumber} of {page_count}"
        self.drawRightString(558, 32, page_str)
        self.restoreState()


def get_custom_stylesheet():
    styles = getSampleStyleSheet()
    
    # Primary Palette
    PRIMARY = colors.HexColor("#1e293b")
    SECONDARY = colors.HexColor("#0284c7")
    TEXT_DARK = colors.HexColor("#0f172a")
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=TEXT_DARK,
        spaceAfter=6
    )
    
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=22,
        leading=26,
        textColor=PRIMARY,
        alignment=0,
        spaceAfter=12
    )
    
    h1_style = ParagraphStyle(
        'CustomH1',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=15,
        leading=19,
        textColor=PRIMARY,
        spaceBefore=14,
        spaceAfter=8,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        'CustomH2',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=12,
        leading=16,
        textColor=SECONDARY,
        spaceBefore=10,
        spaceAfter=6,
        keepWithNext=True
    )
    
    code_style = ParagraphStyle(
        'CustomCode',
        parent=styles['Normal'],
        fontName='Courier',
        fontSize=8.5,
        leading=11,
        textColor=colors.HexColor("#0f172a"),
        backColor=colors.HexColor("#f1f5f9"),
        borderColor=colors.HexColor("#cbd5e1"),
        borderWidth=0.5,
        borderPadding=6,
        spaceBefore=6,
        spaceAfter=6
    )
    
    bullet_style = ParagraphStyle(
        'CustomBullet',
        parent=body_style,
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=4
    )
    
    return {
        'Body': body_style,
        'Title': title_style,
        'H1': h1_style,
        'H2': h2_style,
        'Code': code_style,
        'Bullet': bullet_style
    }


def parse_md_to_flowables(md_path, styles):
    if not os.path.exists(md_path):
        return [Paragraph(f"File not found: {md_path}", styles['Body'])]

    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    flowables = []
    in_code_block = False
    code_lines = []
    in_table = False
    table_data = []

    def clean_text(t):
        # Convert MD bold/italic/links for reportlab Paragraph
        t = t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        t = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', t)
        t = re.sub(r'\*(.*?)\*', r'<i>\1</i>', t)
        t = re.sub(r'`(.*?)`', r'<font face="Courier" color="#0369a1">\1</font>', t)
        t = re.sub(r'\[(.*?)\]\((.*?)\)', r'<font color="#0284c7"><u>\1</u></font>', t)
        return t

    for line in lines:
        raw_line = line.rstrip('\n')
        
        # Code block handling
        if raw_line.startswith('```'):
            if in_code_block:
                code_text = "<br/>".join([clean_text(cl) for cl in code_lines])
                flowables.append(Paragraph(code_text, styles['Code']))
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_lines.append(raw_line)
            continue

        # Table handling
        if '|' in raw_line and len(raw_line.strip()) > 3:
            if '---' in raw_line:
                continue # Skip divider row
            cols = [clean_text(c.strip()) for c in raw_line.strip('|').split('|')]
            table_data.append(cols)
            in_table = True
            continue
        elif in_table:
            if table_data:
                # Render table
                formatted_data = []
                for row_idx, row in enumerate(table_data):
                    formatted_row = []
                    for cell in row:
                        style = ParagraphStyle('TCell', parent=styles['Body'], fontSize=8.5, leading=11)
                        if row_idx == 0:
                            style.fontName = 'Helvetica-Bold'
                            style.textColor = colors.white
                        formatted_row.append(Paragraph(cell, style))
                    formatted_data.append(formatted_row)
                
                col_widths = [504 / max(len(row) for row in table_data)] * max(len(row) for row in table_data)
                t = Table(formatted_data, colWidths=col_widths)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1e293b")),
                    ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 4),
                    ('TOPPADDING', (0,0), (-1,-1), 4),
                    ('LEFTPADDING', (0,0), (-1,-1), 4),
                    ('RIGHTPADDING', (0,0), (-1,-1), 4),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#cbd5e1")),
                    ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor("#f8fafc")])
                ]))
                flowables.append(Spacer(1, 4))
                flowables.append(t)
                flowables.append(Spacer(1, 6))
            table_data = []
            in_table = False

        stripped = raw_line.strip()
        if not stripped:
            continue

        if stripped.startswith('# '):
            flowables.append(Paragraph(clean_text(stripped[2:]), styles['Title']))
        elif stripped.startswith('## '):
            flowables.append(Paragraph(clean_text(stripped[3:]), styles['H1']))
        elif stripped.startswith('### '):
            flowables.append(Paragraph(clean_text(stripped[4:]), styles['H2']))
        elif stripped.startswith('- ') or stripped.startswith('* '):
            flowables.append(Paragraph(f"• {clean_text(stripped[2:])}", styles['Bullet']))
        elif re.match(r'^\d+\.\s', stripped):
            flowables.append(Paragraph(clean_text(stripped), styles['Bullet']))
        elif stripped.startswith('---'):
            flowables.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#e2e8f0"), spaceBefore=8, spaceAfter=8))
        else:
            flowables.append(Paragraph(clean_text(stripped), styles['Body']))

    return flowables


def create_pdf(output_filename, title, md_files):
    doc = SimpleDocTemplate(
        output_filename,
        pagesize=letter,
        leftMargin=54,
        rightMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    styles = get_custom_stylesheet()
    story = []
    
    for idx, md_file in enumerate(md_files):
        if idx > 0:
            story.append(PageBreak())
        elements = parse_md_to_flowables(md_file, styles)
        story.extend(elements)

    doc.build(story, canvasmaker=NumberedCanvas)
    print(f"Generated PDF: {output_filename}")


def create_presentation_pptx(output_filename):
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    NAVY = RGBColor(30, 41, 59)
    AMBER = RGBColor(217, 119, 6)
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(15, 23, 42)
    GRAY_TEXT = RGBColor(100, 116, 139)
    LIGHT_BG = RGBColor(248, 250, 252)
    CARD_BG = RGBColor(241, 245, 249)

    def add_header(slide, title_text, category_text="COGNILOOP RESEARCH"):
        # Header background banner
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        
        # Category label
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = AMBER

        # Main Slide Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

    # --- SLIDE 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()

    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.5))
    tf1 = t_box.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "CogniLoop: Personalizing Open-Domain Video Curricula"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "A Hybrid Framework Combining Thompson Sampling, IRT-Grounded BKT, and Micro-Patterns"
    p1_sub.font.size = Pt(18)
    p1_sub.font.color.rgb = AMBER
    p1_sub.space_before = Pt(15)

    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(10), Inches(1.5))
    tf_m = meta_box.text_frame
    pm1 = tf_m.paragraphs[0]
    pm1.text = "Author & Research Lead: Kavya Aggarwal"
    pm1.font.size = Pt(14)
    pm1.font.color.rgb = WHITE
    pm2 = tf_m.add_paragraph()
    pm2.text = "IEEE Research Project Submission | Release v1.0 | 162/162 Tests Verified"
    pm2.font.size = Pt(12)
    pm2.font.color.rgb = GRAY_TEXT
    pm2.space_before = Pt(5)

    # --- SLIDE 2: Problem & Vision ---
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "The Core Problem & Research Vision")

    cards = [
        ("Passive Consumption Bottleneck", "Modern MOOC video watching is passive, causing rapid cognitive decay and a false sense of mastery without interactive checkpoints.", NAVY),
        ("Content Authoring Bottleneck", "Traditional Intelligent Tutoring Systems (ITS) require hundreds of manual expert hours per video hour, preventing open-domain scaling.", NAVY),
        ("LLM Hallucinations & Blindness", "Standard AI chatbots hallucinate out-of-context facts and lack formal Bayesian cognitive tracking over student mastery trajectories.", NAVY)
    ]

    for idx, (title, desc, color) in enumerate(cards):
        left = Inches(0.8 + idx * 4.0)
        top = Inches(1.8)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = RGBColor(203, 213, 225)
        
        tb = slide2.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), Inches(3.3), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(14)

    # --- SLIDE 3: Four Technical Pillars ---
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "The Four Technical Pillars of CogniLoop")

    pillars = [
        ("1. IRT-Grounded BKT", "Models latent conceptual mastery P(L_n) with dynamic Guess/Slip parameters mapped across Easy (0.30/0.05), Medium (0.20/0.10), and Hard (0.10/0.15) tiers."),
        ("2. Thompson Sampling Bandit", "Solves exploration vs exploitation via Beta-Bernoulli conjugate distributions Beta(alpha, beta) to select optimal difficulty routing vectors."),
        ("3. K-Means Behavioral Pacing", "Classifies video interaction telemetry (speed, pauses, skips, rewinds) into Fast, Steady, and Detail-Oriented student profiles."),
        ("4. Zero-Hallucination RAG", "Indexes lecture transcripts into ChromaDB vector embeddings (all-MiniLM-L6-v2) to ground LLM prompt assembly in exact video context.")
    ]

    for idx, (title, desc) in enumerate(pillars):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.6 + row * 2.7)
        
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = RGBColor(203, 213, 225)

        tb = slide3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = AMBER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(8)

    # --- SLIDE 4: Architecture & Fallback Policy ---
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3-Tier Fault-Tolerant Architecture & Reliability")

    tiers = [
        ("Tier 1: Cloud Groq API", "llama-3.1-8b-instant (~2s response time)", "Primary cloud generation gateway enforcing structured JSON output schemas."),
        ("Tier 2: Local Edge Ollama", "llama3.2 local model (Unlimited offline)", "Automatic edge failover during cloud network isolation or rate limiting."),
        ("Tier 3: Static Curriculum Pool", "Curated Domain Question Bank", "Deterministic safety net guaranteeing continuous examination availability.")
    ]

    for idx, (t_name, t_sub, t_desc) in enumerate(tiers):
        left = Inches(0.8 + idx * 4.0)
        top = Inches(1.8)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = RGBColor(203, 213, 225)

        tb = slide4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), Inches(3.3), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t_name
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p_sub = tf.add_paragraph()
        p_sub.text = t_sub
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = AMBER
        p_sub.space_before = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = t_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(12)

    # --- SLIDE 5: Empirical Validation ---
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Empirical Validation Results (N=60 Simulated Cohort)")

    # Results Table Box
    tb_res = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
    tf_res = tb_res.text_frame
    tf_res.word_wrap = True

    p = tf_res.paragraphs[0]
    p.text = "Key Empirical Outcomes & Statistical Significance:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY

    metrics_text = [
        "• Experimental Group Mean NLG: 0.934 ± 0.192 (Post-Test: 96.0% vs Pre-Test: 37.3%)",
        "• Control Group Mean NLG: 0.363 ± 0.704 (Post-Test: 62.5% vs Pre-Test: 38.0%)",
        "• Gain Delta: +0.571 (+157.3% relative improvement over static control)",
        "• Two-Sample t-Statistic: t = 4.215",
        "• P-Value: p = 0.000088 < 0.0001 (Statistically Significant at alpha = 0.05)",
        "• Automated Test Suite Verification: 162 / 162 Unit & Integration Tests Passed (100% Pass Rate)"
    ]

    for m in metrics_text:
        pm = tf_res.add_paragraph()
        pm.text = m
        pm.font.size = Pt(14)
        pm.font.color.rgb = DARK_TEXT
        pm.space_before = Pt(10)

    prs.save(output_filename)
    print(f"Generated Presentation PPTX: {output_filename}")


def main():
    print("Generating official submission package in 'Submission/' directory...")

    # 1. Individual PDF documents
    create_pdf("Submission/IEEE_RESEARCH_PAPER.pdf", "IEEE Research Paper", ["IEEE_RESEARCH_PAPER.md"])
    create_pdf("Submission/Architecture.pdf", "System Architecture", ["ARCHITECTURE.md"])
    create_pdf("Submission/Simulation_Report.pdf", "Simulation Report", ["SIMULATION_60_PROFILES_REPORT.md"])
    create_pdf("Submission/README.pdf", "Project README", ["README.md"])

    # 2. Master Combined PDF Package (CogniLoop_v1.0.pdf)
    master_sources = [
        "IEEE_RESEARCH_PAPER.md",
        "ARCHITECTURE.md",
        "SIMULATION_60_PROFILES_REPORT.md",
        "DEFENSE_AND_PRESENTATION_GUIDE.md",
        "README.md"
    ]
    create_pdf("Submission/CogniLoop_v1.0.pdf", "CogniLoop Master Submission Package", master_sources)

    # 3. PowerPoint Presentation (Presentation.pptx)
    create_presentation_pptx("Submission/Presentation.pptx")

    print("\nSubmission Package Complete! Files created in 'Submission/':")
    for fname in os.listdir("Submission"):
        fpath = os.path.join("Submission", fname)
        size_kb = os.path.getsize(fpath) / 1024
        print(f" - {fpath} ({size_kb:.1f} KB)")

if __name__ == "__main__":
    main()
